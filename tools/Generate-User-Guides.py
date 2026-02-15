import re
from datetime import date
from math import ceil
from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import importlib


ROOT = Path(__file__).resolve().parents[1]
HELP_DIR = ROOT / "help"
README_PATH = ROOT / "README.md"
TOOL_MAIN_PATH = ROOT / "src" / "SSB_RetuningAutomations.py"
DOCX_TEMPLATE_CANDIDATES = [
    ROOT / "assets" / "templates_docx" / "UserGuideTemplate.docx",
]
PPTX_TEMPLATE_CANDIDATES = [
    ROOT / "assets" / "templates_pptx" / "UserGuideTemplate.pptx",
]
GUIDE_PREFIX = "User-Guide-SSB-Retuning-Automations-v"


def get_tool_version() -> str:
    text = TOOL_MAIN_PATH.read_text(encoding="utf-8")
    match = re.search(r'^TOOL_VERSION\s*=\s*"([0-9]+\.[0-9]+\.[0-9]+)"', text, flags=re.MULTILINE)
    if not match:
        raise RuntimeError("Unable to find TOOL_VERSION in src/SSB_RetuningAutomations.py")
    return match.group(1)


def update_readme_links(version: str) -> None:
    md_name = f"User-Guide-SSB-Retuning-Automations-v{version}.md"
    docx_name = f"User-Guide-SSB-Retuning-Automations-v{version}.docx"
    pptx_name = f"User-Guide-SSB-Retuning-Automations-v{version}.pptx"
    pdf_name = f"User-Guide-SSB-Retuning-Automations-v{version}.pdf"

    new_block = (
        "## 📙 Technical User Guide\n\n"
        "You can find the technical user guide in these formats:\n"
        f"- [Markdown](help/{md_name})\n"
        f"- [Word](help/{docx_name}?raw=true)\n"
        f"- [PowerPoint](help/{pptx_name}?raw=true)\n"
        f"- [PDF](help/{pdf_name}?raw=true)\n"
    )

    readme = README_PATH.read_text(encoding="utf-8")
    pattern = r"## 📙 Technical User Guide\n[\s\S]*?\n---"
    replacement = f"{new_block}\n---"
    if not re.search(pattern, readme):
        raise RuntimeError("Unable to find 'Technical User Guide' section in README.md")

    updated = re.sub(pattern, replacement, readme, count=1)
    README_PATH.write_text(updated, encoding="utf-8")


def guide_paths(version: str) -> dict[str, Path]:
    base_name = f"{GUIDE_PREFIX}{version}"
    return {
        "md": HELP_DIR / f"{base_name}.md",
        "docx": HELP_DIR / f"{base_name}.docx",
        "pptx": HELP_DIR / f"{base_name}.pptx",
        "pdf": HELP_DIR / f"{base_name}.pdf",
    }


def list_versioned_guides(extension: str) -> list[Path]:
    pattern = f"{GUIDE_PREFIX}*.{extension}"
    return sorted((p for p in HELP_DIR.glob(pattern) if p.is_file()), key=lambda p: p.stat().st_mtime, reverse=True)


def pick_source_markdown(target_md: Path) -> Path:
    if target_md.exists():
        return target_md
    candidates = list_versioned_guides("md")
    if not candidates:
        raise FileNotFoundError(f"No markdown guide found in {HELP_DIR}")
    return candidates[0]


def align_help_guides_to_version(version: str) -> dict[str, Path]:
    HELP_DIR.mkdir(parents=True, exist_ok=True)
    paths = guide_paths(version)

    source_md = pick_source_markdown(paths["md"])
    if source_md != paths["md"]:
        paths["md"].unlink(missing_ok=True)
        source_md.rename(paths["md"])

    for ext in ("docx", "pptx", "pdf"):
        target = paths[ext]
        if target.exists():
            continue
        candidates = list_versioned_guides(ext)
        if candidates:
            candidates[0].rename(target)

    return paths


def cleanup_old_versioned_guides(paths: dict[str, Path]) -> None:
    for ext in ("md", "docx", "pptx", "pdf"):
        for candidate in list_versioned_guides(ext):
            if candidate != paths[ext]:
                candidate.unlink(missing_ok=True)


def first_existing_path(candidates: list[Path]) -> Path | None:
    for path in candidates:
        if path.exists():
            return path
    return None


def normalize_markdown_inline(text: str) -> str:
    normalized = text.replace("`", '"').replace("'", '"')
    normalized = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1 (\2)", normalized)
    return normalized


def clean_heading_text(text: str) -> str:
    """Remove manual markdown numbering so templates can apply automatic numbering."""
    cleaned = text.strip()

    # e.g. "1) Title", "2.3 Subtitle", "4.1.2 Topic"
    cleaned = re.sub(r"^\d+(?:[.)]|(?:\.\d+)*\.?)[\s]+", "", cleaned)

    # e.g. "A) Title", "B) Title", "A.B) Title", "a) Title"
    cleaned = re.sub(r"^[A-Za-z](?:\.[A-Za-z])*\)\s+", "", cleaned)

    return cleaned.strip()


def parse_table_row(line: str) -> list[str]:
    stripped = line.strip().strip("|")
    return [cell.strip() for cell in stripped.split("|")]


def is_table_separator(line: str) -> bool:
    stripped = line.strip()
    return bool(re.fullmatch(r"\|?[\s:-]+(\|[\s:-]+)+\|?", stripped))


def build_docx_from_markdown(md_file: Path, docx_file: Path) -> None:
    # ------------------- Word-only helpers (subfunctions) ------------------- #
    def markdown_segments(text: str) -> list[tuple[str, bool]]:
        """Return [(segment, is_bold)] for markdown strings with **bold** markers."""
        clean = normalize_markdown_inline(text)
        parts = re.split(r"(\*\*.*?\*\*)", clean)
        segments: list[tuple[str, bool]] = []
        for part in parts:
            if not part:
                continue
            if part.startswith("**") and part.endswith("**") and len(part) > 4:
                segments.append((part[2:-2], True))
            else:
                segments.append((part, False))
        return segments

    def add_markdown_paragraph(doc: Document, text: str, style: str | None = None) -> None:
        paragraph = doc.add_paragraph(style=style) if style else doc.add_paragraph()
        for segment, is_bold in markdown_segments(text):
            run = paragraph.add_run(segment)
            run.bold = is_bold

    def add_plain_paragraph(doc: Document, text: str, style: str | None = None) -> None:
        paragraph = doc.add_paragraph(style=style) if style else doc.add_paragraph()
        paragraph.add_run(normalize_markdown_inline(text))

    def enable_toc_update_on_open(doc: Document) -> None:
        settings = doc.settings.element
        node = settings.find(qn("w:updateFields"))
        if node is None:
            node = OxmlElement("w:updateFields")
            settings.append(node)
        node.set(qn("w:val"), "true")

    def update_header_date(doc: Document) -> None:
        today = date.today().isoformat()

        def local_name(tag: str) -> str:
            return tag.split("}")[-1]

        for section in doc.sections:
            header_el = section.header._element
            for sdt in header_el.iter():
                if local_name(sdt.tag) != "sdt":
                    continue

                is_date_alias = False
                for node in sdt.iter():
                    if local_name(node.tag) == "alias" and node.get(qn("w:val")) == "Date":
                        is_date_alias = True
                        break

                if not is_date_alias:
                    continue

                for node in sdt.iter():
                    if local_name(node.tag) == "t":
                        node.text = today
                        break

    def find_template_anchor(doc: Document):
        for paragraph in doc.paragraphs:
            if paragraph.text.strip().lower() == "(autogenerated-content-start)":
                return paragraph
        return None

    def remove_from_anchor_to_end(doc: Document, anchor_paragraph) -> None:
        body = doc._element.body
        start = body.index(anchor_paragraph._p)
        for idx in range(len(body) - 1, start - 1, -1):
            body.remove(body[idx])

    # ------------------------------ Word logic ------------------------------ #
    text = md_file.read_text(encoding="utf-8")
    lines = text.splitlines()

    template_path = first_existing_path(DOCX_TEMPLATE_CANDIDATES)
    doc = Document(str(template_path)) if template_path else Document()

    anchor = find_template_anchor(doc)
    if anchor is not None:
        remove_from_anchor_to_end(doc, anchor)

    seen_heading1 = False  # Track if we've already written the first Heading 1 (## ...)

    i = 0
    while i < len(lines):
        line = lines[i]
        s = line.rstrip()

        if not s:
            # Skip blank lines around headings and around horizontal rules (---)
            def is_heading(line_: str) -> bool:
                t = (line_ or "").strip()
                return t.startswith("# ") or t.startswith("## ") or t.startswith("### ") or t.startswith("#### ")

            def is_hr(line_: str) -> bool:
                return (line_ or "").strip() == "---"

            # Previous non-empty line
            j = i - 1
            while j >= 0 and not lines[j].strip():
                j -= 1
            prev = lines[j] if j >= 0 else ""

            # Next non-empty line
            k = i + 1
            while k < len(lines) and not lines[k].strip():
                k += 1
            nxt = lines[k] if k < len(lines) else ""

            # If blank line is before/after a heading or an '---', drop it
            if is_heading(prev) or is_heading(nxt) or is_hr(prev) or is_hr(nxt):
                i += 1
                continue

            doc.add_paragraph("")
            i += 1
            continue

        if s == "---":
            i += 1
            # Also skip any blank lines immediately after the horizontal rule
            while i < len(lines) and not lines[i].strip():
                i += 1
            continue

        if s.startswith("| ") and s.endswith(" |") and i + 1 < len(lines) and is_table_separator(lines[i + 1]):
            rows: list[list[str]] = [parse_table_row(s)]
            i += 2
            while i < len(lines):
                candidate = lines[i].rstrip()
                if not (candidate.startswith("| ") and candidate.endswith(" |")):
                    break
                rows.append(parse_table_row(candidate))
                i += 1

            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = "Table Grid"
                for row_idx, row in enumerate(rows):
                    for col_idx, cell in enumerate(row):
                        paragraph = table.cell(row_idx, col_idx).paragraphs[0]
                        for seg, is_bold in markdown_segments(cell):
                            run = paragraph.add_run(seg)
                            run.bold = is_bold or row_idx == 0
            continue

        if s.startswith("# "):
            # Markdown H1 is the document title (not part of heading numbering hierarchy).
            add_plain_paragraph(doc, clean_heading_text(s[2:].strip()), style="Title")

        elif s.startswith("## "):
            # Each Heading 1 starts on a new page, except the first one (usually right after the Title)
            if seen_heading1:
                doc.add_page_break()
            add_plain_paragraph(doc, clean_heading_text(s[3:].strip()), style="Heading 1")
            seen_heading1 = True

        elif s.startswith("### "):
            add_plain_paragraph(doc, clean_heading_text(s[4:].strip()), style="Heading 2")

        elif s.startswith("#### "):
            add_plain_paragraph(doc, clean_heading_text(s[5:].strip()), style="Heading 3")

        # Bullets (support sub-bullets by indentation)
        elif re.match(r"^\s*-\s+", s):
            m = re.match(r"^(\s*)-\s+(.*)$", s)
            indent = len(m.group(1).replace("\t", "    "))
            item_text = m.group(2).strip()

            # 0 spaces => bullet normal; >=2 spaces => sub-bullet
            style = "List Bullet 2" if indent >= 2 else "List Bullet"
            add_markdown_paragraph(doc, item_text, style=style)

        elif re.match(r"^[A-Za-z](?:\.[A-Za-z])*\)\s+", s):
            # Examples: A) , B) , a) , A.B) , a.b)
            add_markdown_paragraph(doc, re.sub(r"^[A-Za-z](?:\.[A-Za-z])*\)\s+", "", s), style="List Number")

        elif re.match(r"^\d+\.\s+", s):
            # Keep Markdown numbering to avoid Word continuing previous lists
            add_markdown_paragraph(doc, s.strip(), style="List Paragraph")

        else:
            add_markdown_paragraph(doc, s)

        i += 1

    enable_toc_update_on_open(doc)
    update_header_date(doc)
    doc.save(docx_file)


def build_pptx_summary(md_file: Path, pptx_file: Path, version: str) -> None:
    # ------------------- PPT-only helpers (subfunctions) ------------------- #
    def _strip_bold(text: str) -> str:
        return re.sub(r"\*\*(.*?)\*\*", r"\1", text)

    def _norm(text: str) -> str:
        return normalize_markdown_inline(_strip_bold(text))

    def _indent_to_level(indent_spaces: int) -> int:
        # 0-1 => 0, 2-7 => 1, 8-11 => 2, ...
        if indent_spaces < 2:
            return 0
        lvl = indent_spaces // 4
        if lvl == 0:
            lvl = 1
        return min(lvl, 4)

    def _parse_sections_for_ppt(md_path: Path) -> list[dict]:
        """
        Parse markdown into sections (##) and blocks.
        - ### becomes block type 'h3' (we'll use it as slide subtitle)
        - #### becomes block type 'h4'
        - bullets preserve indentation
        - supports numeric ordered (1.) and alpha ordered (A), A.B))
        """
        lines = md_path.read_text(encoding="utf-8").splitlines()
        sections: list[dict] = []
        current = {"title": "Overview", "blocks": []}

        bullet_pat = re.compile(r"^(\s*)-\s+(.*)$")
        num_pat = re.compile(r"^(\s*)(\d+)\.\s+(.*)$")
        alpha_pat = re.compile(r"^(\s*)([A-Za-z](?:\.[A-Za-z])*)\)\s+(.*)$")

        i = 0
        while i < len(lines):
            raw = lines[i].rstrip()
            line = raw.strip()

            if not line or line == "---":
                i += 1
                continue

            if line.startswith("# "):
                i += 1
                continue

            if line.startswith("## "):
                if current["blocks"]:
                    sections.append(current)
                current = {"title": clean_heading_text(line[3:].strip()), "blocks": []}
                i += 1
                continue

            if line.startswith("### "):
                current["blocks"].append({"type": "h3", "text": clean_heading_text(line[4:].strip())})
                i += 1
                continue

            if line.startswith("#### "):
                current["blocks"].append({"type": "h4", "text": clean_heading_text(line[5:].strip())})
                i += 1
                continue

            if line.startswith("| ") and line.endswith(" |") and i + 1 < len(lines) and is_table_separator(lines[i + 1]):
                rows = [parse_table_row(line)]
                i += 2
                while i < len(lines):
                    cand_raw = lines[i].rstrip()
                    cand = cand_raw.strip()
                    if not (cand.startswith("| ") and cand.endswith(" |")):
                        break
                    rows.append(parse_table_row(cand))
                    i += 1
                current["blocks"].append({"type": "table", "rows": rows})
                continue

            mb = bullet_pat.match(raw)
            mn = num_pat.match(raw)
            ma = alpha_pat.match(raw)

            if mb or mn or ma:
                items: list[dict] = []
                mode = "bullet" if mb else ("num" if mn else "alpha")

                while i < len(lines):
                    raw_c = lines[i].rstrip("\n")
                    c = raw_c.strip()

                    if not c or c == "---":
                        break
                    if c.startswith("## ") or c.startswith("### ") or c.startswith("#### "):
                        break
                    if c.startswith("| "):
                        break

                    mb2 = bullet_pat.match(raw_c)
                    mn2 = num_pat.match(raw_c)
                    ma2 = alpha_pat.match(raw_c)

                    if mode == "bullet" and mb2:
                        indent_spaces = len(mb2.group(1).replace("\t", "    "))
                        items.append({
                            "text": _norm(mb2.group(2).strip()),
                            "level": _indent_to_level(indent_spaces),
                            "marker": None,
                        })
                        i += 1
                        continue

                    if mode == "num" and mn2:
                        indent_spaces = len(mn2.group(1).replace("\t", "    "))
                        marker = f"{mn2.group(2)}."
                        items.append({
                            "text": _norm(mn2.group(3).strip()),
                            "level": _indent_to_level(indent_spaces),
                            "marker": marker,
                        })
                        i += 1
                        continue

                    if mode == "alpha" and ma2:
                        indent_spaces = len(ma2.group(1).replace("\t", "    "))
                        marker = f"{ma2.group(2)})"
                        items.append({
                            "text": _norm(ma2.group(3).strip()),
                            "level": _indent_to_level(indent_spaces),
                            "marker": marker,
                        })
                        i += 1
                        continue

                    break

                current["blocks"].append({
                    "type": "list",
                    "ordered": mode in ("num", "alpha"),
                    "items": items,
                })
                continue

            current["blocks"].append({"type": "paragraph", "text": _norm(line)})
            i += 1

        if current["blocks"]:
            sections.append(current)

        return sections

    def _block_units(block: dict) -> int:
        # Rough estimator to fit content in a slide.
        if block["type"] == "h4":
            return 2
        if block["type"] == "paragraph":
            return max(1, ceil(len(block["text"]) / 95))
        if block["type"] == "list":
            # 1 for list header overhead + per item length
            return 1 + sum(max(1, ceil(len(it["text"]) / 95)) for it in block["items"])
        if block["type"] == "table":
            return 6 + len(block["rows"])
        return 2

    def _style_table(table) -> None:
        header_fill = RGBColor(31, 78, 121)
        band_fill = RGBColor(242, 246, 252)
        text_dark = RGBColor(40, 40, 40)

        for col in range(len(table.columns)):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(14)

        for row in range(1, len(table.rows)):
            for col in range(len(table.columns)):
                cell = table.cell(row, col)
                if row % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = band_fill
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = text_dark
                p.font.size = Pt(12)

    def _add_table_slide(prs: Presentation, slide_title: str, subtitle: str | None, rows: list[list[str]], continuation: bool = False) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        # Subtitle on table slides too (small)
        if subtitle:
            p0 = tf.paragraphs[0]
            p0.text = subtitle if not continuation else f"{subtitle} (cont.)"
            p0.level = 0
            p0.font.bold = True
            p0.font.size = Pt(18)
        else:
            # ensure we have a first paragraph ready even if no subtitle
            tf.paragraphs[0].text = ""

        content = slide.shapes.placeholders[1]
        x, y, w, h = content.left, content.top, content.width, content.height
        tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
        table = tbl_shape.table

        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                table.cell(r, c).text = _norm(value)

        _style_table(table)

    def _render_content(tf, blocks: list[dict]) -> None:
        from pptx.oxml.ns import qn as pptx_qn
        from pptx.oxml.xmlchemy import OxmlElement as PptxOxmlElement

        def _set_bullets(p, enabled: bool) -> None:
            pPr = p._p.get_or_add_pPr()
            for tag in ("a:buAutoNum", "a:buChar", "a:buBlip", "a:buFont", "a:buNone"):
                el = pPr.find(pptx_qn(tag))
                if el is not None:
                    pPr.remove(el)
            if enabled:
                bu = PptxOxmlElement("a:buChar")
                bu.set("char", "•")
                pPr.insert(0, bu)
            else:
                pPr.insert(0, PptxOxmlElement("a:buNone"))

        tf.clear()
        first = True

        for block in blocks:
            if block["type"] == "h4":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = block["text"]
                p.level = 0
                _set_bullets(p, enabled=False)
                p.font.bold = True
                p.font.size = Pt(16)
                continue

            if block["type"] == "paragraph":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = block["text"]
                p.level = 0
                _set_bullets(p, enabled=False)
                p.font.size = Pt(15)
                continue

            if block["type"] == "list":
                for idx, it in enumerate(block["items"], start=1):
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False

                    nesting = int(it.get("level", 0))
                    p.level = 1 + nesting
                    p.font.size = Pt(15)

                    if block["ordered"]:
                        marker = it.get("marker") or f"{idx}."
                        p.text = f"{marker} {it['text']}"
                        _set_bullets(p, enabled=False)
                    else:
                        p.text = it["text"]
                        _set_bullets(p, enabled=True)

    # ------------------------------ PPT logic ------------------------------ #
    sections = _parse_sections_for_ppt(md_file)

    template_path = first_existing_path(PPTX_TEMPLATE_CANDIDATES)
    prs = Presentation(str(template_path)) if template_path else Presentation()

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Technical User Guide Summary"
    slide.placeholders[1].text = f"SSB Retuning Automations (v{version})"

    max_units = 20

    for section in sections:
        slide_title = section["title"]
        blocks = section["blocks"]

        # Group blocks by ### (h3). Each group becomes one or more slides.
        groups: list[dict] = []
        current_subtitle: str | None = None
        current_blocks: list[dict] = []

        for b in blocks:
            if b["type"] == "h3":
                # flush previous group
                if current_subtitle is not None or current_blocks:
                    groups.append({"subtitle": current_subtitle, "blocks": current_blocks})
                current_subtitle = b["text"]
                current_blocks = []
            else:
                current_blocks.append(b)

        if current_subtitle is not None or current_blocks:
            groups.append({"subtitle": current_subtitle, "blocks": current_blocks})

        for g in groups:
            subtitle = g["subtitle"]
            gblocks = g["blocks"]

            # paginate blocks for this group
            buffer: list[dict] = []
            used = 0
            continuation = False

            def flush() -> None:
                nonlocal buffer, used, continuation
                slide_local = prs.slides.add_slide(prs.slide_layouts[1])

                # Title = ## (slide title) + ### (subtitle)
                title_shape = slide_local.shapes.title
                title_shape.text = slide_title
                if subtitle:
                    p = title_shape.text_frame.add_paragraph()
                    p.text = subtitle if not continuation else f"{subtitle} (cont.)"
                    p.level = 0
                    p.font.bold = False
                    p.font.size = Pt(20)
                    p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

                # Body content (NO subtitle here)
                _render_content(slide_local.shapes.placeholders[1].text_frame, buffer)

                buffer = []
                used = 0
                continuation = True

            i = 0
            while i < len(gblocks):
                b = gblocks[i]

                if b["type"] == "table":
                    if buffer:
                        flush()
                    _add_table_slide(prs, slide_title, subtitle, b["rows"], continuation=continuation)
                    continuation = True
                    i += 1
                    continue

                u = _block_units(b)
                if buffer and used + u > max_units:
                    flush()

                buffer.append(b)
                used += u
                i += 1

            if buffer:
                flush()

    prs.save(pptx_file)


def build_pdf_from_markdown(md_file: Path, pdf_file: Path) -> None:
    if importlib.util.find_spec("reportlab") is None:
        raise RuntimeError("Missing dependency 'reportlab'. Install requirements.txt to generate PDF guides.")

    colors = importlib.import_module("reportlab.lib.colors")
    pagesizes = importlib.import_module("reportlab.lib.pagesizes")
    styles_mod = importlib.import_module("reportlab.lib.styles")
    units = importlib.import_module("reportlab.lib.units")
    platypus = importlib.import_module("reportlab.platypus")

    A4 = pagesizes.A4
    ParagraphStyle = styles_mod.ParagraphStyle
    getSampleStyleSheet = styles_mod.getSampleStyleSheet
    cm = units.cm
    Paragraph = platypus.Paragraph
    SimpleDocTemplate = platypus.SimpleDocTemplate
    Spacer = platypus.Spacer
    Table = platypus.Table
    TableStyle = platypus.TableStyle

    styles = getSampleStyleSheet()
    body_style = ParagraphStyle("Body", parent=styles["BodyText"], leading=14, spaceAfter=6)
    h1_style = ParagraphStyle("H1", parent=styles["Heading1"], spaceBefore=8, spaceAfter=6)
    h2_style = ParagraphStyle("H2", parent=styles["Heading2"], spaceBefore=6, spaceAfter=4)
    h3_style = ParagraphStyle("H3", parent=styles["Heading3"], spaceBefore=4, spaceAfter=3)

    def inline_html(text: str) -> str:
        text = normalize_markdown_inline(text)
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        return re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)

    doc = SimpleDocTemplate(str(pdf_file), pagesize=A4, leftMargin=1.8 * cm, rightMargin=1.8 * cm, topMargin=1.8 * cm, bottomMargin=1.8 * cm)
    story = []
    lines = md_file.read_text(encoding="utf-8").splitlines()

    i = 0
    while i < len(lines):
        raw = lines[i].rstrip()
        line = raw.strip()

        if not line or line == "---":
            story.append(Spacer(1, 0.2 * cm))
            i += 1
            continue

        if line.startswith("| ") and line.endswith(" |") and i + 1 < len(lines) and is_table_separator(lines[i + 1]):
            rows = [parse_table_row(line)]
            i += 2
            while i < len(lines):
                candidate = lines[i].strip()
                if not (candidate.startswith("| ") and candidate.endswith(" |")):
                    break
                rows.append(parse_table_row(candidate))
                i += 1
            table_data = [[Paragraph(inline_html(cell), body_style) for cell in row] for row in rows]
            table = Table(table_data, repeatRows=1)
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e2e8f0")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(table)
            story.append(Spacer(1, 0.2 * cm))
            continue

        if line.startswith("# "):
            story.append(Paragraph(inline_html(clean_heading_text(line[2:].strip())), h1_style))
        elif line.startswith("## "):
            story.append(Paragraph(inline_html(clean_heading_text(line[3:].strip())), h1_style))
        elif line.startswith("### "):
            story.append(Paragraph(inline_html(clean_heading_text(line[4:].strip())), h2_style))
        elif line.startswith("#### "):
            story.append(Paragraph(inline_html(clean_heading_text(line[5:].strip())), h3_style))
        elif re.match(r"^\s*-\s+", raw):
            bullet = re.sub(r"^\s*-\s+", "", raw)
            story.append(Paragraph(inline_html(f"• {bullet}"), body_style))
        else:
            story.append(Paragraph(inline_html(line), body_style))

        i += 1

    doc.build(story)



if __name__ == "__main__":
    version = get_tool_version()
    update_readme_links(version)

    paths = align_help_guides_to_version(version)
    build_docx_from_markdown(paths["md"], paths["docx"])
    build_pptx_summary(paths["md"], paths["pptx"], version)
    build_pdf_from_markdown(paths["md"], paths["pdf"])
    cleanup_old_versioned_guides(paths)

    print(f"Tool version detected: v{version}")
    print(f"Using markdown: {paths['md']}")
    print(f"Generated: {paths['docx']}")
    print(f"Generated: {paths['pptx']}")
    print(f"Generated: {paths['pdf']}")
    print("README technical guide links updated.")
    print("Old versioned guides cleaned from help/.")
