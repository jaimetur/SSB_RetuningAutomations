#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse
import re
from datetime import date
from math import ceil
from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm as DocxCm
from docx.shared import Pt as DocxPt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
import importlib
import os
import shutil
import subprocess

# =========== WORD SETTINGS =============
WORD_TABLE_HEADER_PT = 9            # Font Size for Table Header
WORD_TABLE_BODY_PT = 8              # Font Size for Table Body
WORD_LIST_STYLE = "List Paragraph"  # Style for List Paragraph (Must exists in the template)
# =======================================

# =========== PPT SETTINGS ==============
PPT_FONT_SIZE_H3 = 20               # Define Font Size for Headers 3 (Header 3 are in the Tittle of the slide as subtitle)
PPT_FONT_SIZE_H4 = 14               # Define Font Size for Headers 4 (Header 2 and 3 are in the Tittle of the slide)
PPT_FONT_SIZE_PARAGRAPH = 12        # Font size for normal paragraph
PPT_FONT_SIZE_LISTS = 12            # Font size for lists
PPT_FONT_SIZE_TABLES_HEADER = 10    # Font size for table header
PPT_FONT_SIZE_TABLES_BODY = 7       # Font size for table body

PPT_MAX_CONTENT_LINES = 27          # Max number of lines per slide
PPT_MAX_TABLE_DATA_ROWS = 15        # Max number of Rows per table (header doesn' t count)
PPT_CUT_SLIDE_WEIGHT = 130          # Increasing this value makes to add more content to previous slide
                                    # Decreasing this value makes the lists and paragraph to be cut to next slide earlier

PPT_COMBINE_MAX_DATA_ROWS = 12      # Only try to combine small tables (header excluded)
PPT_GAP = Pt(6)                     # Gap between rows
PPT_MIN_TEXT_H = Pt(70)             # Min Text Height to combine with a table in the same slide
PPT_MIN_TABLE_H = Pt(160)           # Min Table Height to combine with text in the same slide
# =======================================

# ================ PATHS ================
ROOT = Path(__file__).resolve().parents[1]
HELP_DIR = ROOT / "help"
README_PATH = ROOT / "README.md"
GUIDE_PREFIX = "User-Guide-SSB-Retuning-Automations-v"
TOOL_MAIN_PATH = ROOT / "src" / "SSB_RetuningAutomations.py"
DOCX_TEMPLATE_CANDIDATES = [ROOT / "assets" / "templates_docx" / "UserGuideTemplate.docx"]
PPTX_TEMPLATE_CANDIDATES = [ROOT / "assets" / "templates_pptx" / "UserGuideTemplate.pptx"]
# =======================================

IMAGE_MD_PATTERN = re.compile(r'^!\[(?P<alt>[^\]]*)\]\((?P<src>[^)\s]+)(?:\s+"[^"]*")?\)\s*$')


def get_tool_version() -> str:
    text = TOOL_MAIN_PATH.read_text(encoding="utf-8")
    match = re.search(r'^TOOL_VERSION\s*=\s*"([0-9]+\.[0-9]+\.[0-9]+)"', text, flags=re.MULTILINE)
    if not match:
        raise RuntimeError("Unable to find TOOL_VERSION in src/SSB_RetuningAutomations.py")
    return match.group(1)


def update_readme_links(version: str) -> None:
    md_name = f"User-Guide-SSB-Retuning-Automations-v{version}.md"
    docx_pdf_name = f"User-Guide-SSB-Retuning-Automations-v{version}.docx.pdf"
    pptx_pdf_name = f"User-Guide-SSB-Retuning-Automations-v{version}.pptx.pdf"

    new_block = (
        "## 📙 Technical User Guide\n\n"
        "You can find the technical user guide in these formats:\n"
        f"- [Markdown](help/{md_name})\n"
        f"- [Word (PDF)](help/{docx_pdf_name}?raw=true)\n"
        f"- [PowerPoint (PDF)](help/{pptx_pdf_name}?raw=true)\n"
    )

    readme = README_PATH.read_text(encoding="utf-8")
    pattern = r"## 📙 Technical User Guide\n[\s\S]*?\n---"
    replacement = f"{new_block}\n---"
    if not re.search(pattern, readme):
        raise RuntimeError("Unable to find 'Technical User Guide' section in README.md")

    updated = re.sub(pattern, replacement, readme, count=1)
    README_PATH.write_text(updated, encoding="utf-8")


def guide_paths(version: str) -> dict[str, Path]:
    base_name = f"{GUIDE_PREFIX}{version}"
    return {
        "md": HELP_DIR / f"{base_name}.md",
        "docx": HELP_DIR / f"{base_name}.docx",
        "pptx": HELP_DIR / f"{base_name}.pptx",
        "docx_pdf": HELP_DIR / f"{base_name}.docx.pdf",
        "pptx_pdf": HELP_DIR / f"{base_name}.pptx.pdf",
    }


def list_versioned_guides(extension: str) -> list[Path]:
    pattern = f"{GUIDE_PREFIX}*.{extension}"
    return sorted((p for p in HELP_DIR.glob(pattern) if p.is_file()), key=lambda p: p.stat().st_mtime, reverse=True)


def pick_source_markdown(target_md: Path) -> Path:
    if target_md.exists():
        return target_md
    candidates = list_versioned_guides("md")
    if not candidates:
        raise FileNotFoundError(f"No markdown guide found in {HELP_DIR}")
    return candidates[0]


def align_help_guides_to_version(version: str) -> dict[str, Path]:
    HELP_DIR.mkdir(parents=True, exist_ok=True)
    paths = guide_paths(version)

    source_md = pick_source_markdown(paths["md"])
    if source_md != paths["md"]:
        paths["md"].unlink(missing_ok=True)
        source_md.rename(paths["md"])

    legacy_to_target = {
        "docx": paths["docx"],
        "pptx": paths["pptx"],
        "docx.pdf": paths["docx_pdf"],
        "pptx.pdf": paths["pptx_pdf"],
        "pdf": paths["docx_pdf"],
    }
    for ext, target in legacy_to_target.items():
        if target.exists():
            continue
        candidates = list_versioned_guides(ext)
        if candidates:
            candidates[0].rename(target)

    return paths


def cleanup_old_versioned_guides(paths: dict[str, Path]) -> None:
    keep_by_ext = {
        "md": paths["md"],
        "docx": paths["docx"],
        "pptx": paths["pptx"],
        "docx.pdf": paths["docx_pdf"],
        "pptx.pdf": paths["pptx_pdf"],
        "pdf": paths["docx_pdf"],
    }
    for ext, keep_path in keep_by_ext.items():
        for candidate in list_versioned_guides(ext):
            if candidate != keep_path:
                candidate.unlink(missing_ok=True)


def first_existing_path(candidates: list[Path]) -> Path | None:
    for path in candidates:
        if path.exists():
            return path
    return None


def normalize_markdown_inline(text: str) -> str:
    normalized = text.replace("`", '"').replace("'", '"')
    normalized = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1 (\2)", normalized)
    return normalized


def parse_markdown_image_line(line: str) -> tuple[str, str] | None:
    match = IMAGE_MD_PATTERN.match(line.strip())
    if not match:
        return None
    return match.group("alt").strip(), match.group("src").strip()


def is_html_comment_line(line: str) -> bool:
    stripped = line.strip()
    return stripped.startswith("<!--") and stripped.endswith("-->")


def resolve_markdown_image_path(md_file: Path, source: str) -> Path | None:
    clean_source = source.split("?", 1)[0].strip()
    if not clean_source:
        return None
    source_path = Path(clean_source)
    candidate = source_path if source_path.is_absolute() else (md_file.parent / source_path)
    try:
        resolved = candidate.resolve()
    except OSError:
        return None
    return resolved if resolved.exists() and resolved.is_file() else None


def clean_heading_text(text: str) -> str:
    """Remove manual markdown numbering so templates can apply automatic numbering."""
    cleaned = text.strip()

    # e.g. "1) Title", "2.3 Subtitle", "4.1.2 Topic"
    cleaned = re.sub(r"^\d+(?:[.)]|(?:\.\d+)*\.?)[\s]+", "", cleaned)

    # e.g. "A) Title", "B) Title", "A.B) Title", "a) Title"
    cleaned = re.sub(r"^[A-Za-z](?:\.[A-Za-z])*\)\s+", "", cleaned)

    return cleaned.strip()


def parse_table_row(line: str) -> list[str]:
    stripped = line.strip().strip("|")
    return [cell.strip() for cell in stripped.split("|")]


def is_table_separator(line: str) -> bool:
    stripped = line.strip()
    return bool(re.fullmatch(r"\|?[\s:-]+(\|[\s:-]+)+\|?", stripped))


# Update Word Fields (including ToC) and export into PDF
def try_update_docx_fields_and_export_pdf(docx_file: Path, pdf_file: Path) -> bool:
    """Best effort on Windows: update TOC/fields via Word automation and export PDF."""
    if os.name != "nt":
        return False

    print(f"\tUpdating Word fields and exporting PDF from Word...")
    ps_script = f"""
$ErrorActionPreference = 'Stop'
$word = $null
try {{
    $docPath = '{str(docx_file)}'
    $pdfPath = '{str(pdf_file)}'
    $word = New-Object -ComObject Word.Application
    $word.Visible = $false
    $word.DisplayAlerts = 0
    $doc = $word.Documents.Open($docPath)
    $doc.Fields.Update() | Out-Null
    foreach ($toc in $doc.TablesOfContents) {{ $toc.Update() | Out-Null }}
    $doc.Save()
    $wdExportFormatPDF = 17
    $wdExportOptimizeForPrint = 0
    $wdExportAllDocument = 0
    $wdExportDocumentContent = 0
    $wdExportCreateHeadingBookmarks = 1

    $doc.ExportAsFixedFormat(
        $pdfPath,
        $wdExportFormatPDF,
        $false,
        $wdExportOptimizeForPrint,
        $wdExportAllDocument,
        1,
        1,
        $wdExportDocumentContent,
        $false,
        $true,
        $wdExportCreateHeadingBookmarks,
        $true,
        $true,
        $false
    )

    $doc.Close()
    exit 0
}} catch {{
    if ($doc) {{ $doc.Close() }}
    exit 1
}} finally {{
    if ($word) {{ $word.Quit() }}
}}
"""

    result = subprocess.run(
        ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps_script],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="ignore",
    )
    return result.returncode == 0 and pdf_file.exists()



# Only used when we execute this script in a Non-Windows OS or without Microsoft Word installed
# Markdown to PDF Converter
def build_pdf_from_markdown(md_file: Path, pdf_file: Path) -> None:
    if importlib.util.find_spec("reportlab") is None:
        raise RuntimeError("Missing dependency 'reportlab'. Install requirements.txt to generate PDF guides.")

    print(f"\tGenerating PDF from Markdown...")
    colors = importlib.import_module("reportlab.lib.colors")
    pagesizes = importlib.import_module("reportlab.lib.pagesizes")
    styles_mod = importlib.import_module("reportlab.lib.styles")
    units = importlib.import_module("reportlab.lib.units")
    platypus = importlib.import_module("reportlab.platypus")

    A4 = pagesizes.A4
    ParagraphStyle = styles_mod.ParagraphStyle
    getSampleStyleSheet = styles_mod.getSampleStyleSheet
    cm = units.cm
    Paragraph = platypus.Paragraph
    SimpleDocTemplate = platypus.SimpleDocTemplate
    Spacer = platypus.Spacer
    Table = platypus.Table
    TableStyle = platypus.TableStyle

    styles = getSampleStyleSheet()
    body_style = ParagraphStyle("Body", parent=styles["BodyText"], leading=14, spaceAfter=6)
    h1_style = ParagraphStyle("H1", parent=styles["Heading1"], spaceBefore=8, spaceAfter=6)
    h2_style = ParagraphStyle("H2", parent=styles["Heading2"], spaceBefore=6, spaceAfter=4)
    h3_style = ParagraphStyle("H3", parent=styles["Heading3"], spaceBefore=4, spaceAfter=3)

    def inline_html(text: str) -> str:
        text = normalize_markdown_inline(text)
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        return re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)

    doc = SimpleDocTemplate(str(pdf_file), pagesize=A4, leftMargin=1.8 * cm, rightMargin=1.8 * cm, topMargin=1.8 * cm, bottomMargin=1.8 * cm)
    story = []
    lines = md_file.read_text(encoding="utf-8").splitlines()

    i = 0
    while i < len(lines):
        raw = lines[i].rstrip()
        line = raw.strip()

        if not line or line == "---":
            story.append(Spacer(1, 0.2 * cm))
            i += 1
            continue

        if is_html_comment_line(line):
            i += 1
            continue

        if line.startswith("| ") and line.endswith(" |") and i + 1 < len(lines) and is_table_separator(lines[i + 1]):
            rows = [parse_table_row(line)]
            i += 2
            while i < len(lines):
                candidate = lines[i].strip()
                if not (candidate.startswith("| ") and candidate.endswith(" |")):
                    break
                rows.append(parse_table_row(candidate))
                i += 1
            table_data = [[Paragraph(inline_html(cell), body_style) for cell in row] for row in rows]
            table = Table(table_data, repeatRows=1)
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e2e8f0")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(table)
            story.append(Spacer(1, 0.2 * cm))
            continue

        image_info = parse_markdown_image_line(line)
        if image_info:
            _, image_src = image_info
            image_path = resolve_markdown_image_path(md_file, image_src)
            if image_path:
                img = platypus.Image(str(image_path))
                max_width = doc.width
                max_height = 12 * cm
                iw = float(getattr(img, "imageWidth", 0) or 0)
                ih = float(getattr(img, "imageHeight", 0) or 0)
                if iw > 0 and ih > 0:
                    ratio = min(max_width / iw, max_height / ih, 1.0)
                    img.drawWidth = iw * ratio
                    img.drawHeight = ih * ratio
                story.append(img)
                story.append(Spacer(1, 0.2 * cm))
            i += 1
            continue

        if line.startswith("# "):
            story.append(Paragraph(inline_html(clean_heading_text(line[2:].strip())), h1_style))
        elif line.startswith("## "):
            story.append(Paragraph(inline_html(clean_heading_text(line[3:].strip())), h1_style))
        elif line.startswith("### "):
            story.append(Paragraph(inline_html(clean_heading_text(line[4:].strip())), h2_style))
        elif line.startswith("#### "):
            story.append(Paragraph(inline_html(clean_heading_text(line[5:].strip())), h3_style))
        elif re.match(r"^\s*-\s+", raw):
            bullet = re.sub(r"^\s*-\s+", "", raw)
            story.append(Paragraph(inline_html(f"• {bullet}"), body_style))
        else:
            story.append(Paragraph(inline_html(line), body_style))

        i += 1

    doc.build(story)

# Markdown to Word Converter
def build_docx_from_markdown(md_file: Path, docx_file: Path, version: str) -> None:
    print(f"\tGenerating Word from Markdown...")
    # ------------------- Word-only helpers (subfunctions) ------------------- #
    def markdown_segments(text: str) -> list[tuple[str, bool]]:
        """Return [(segment, is_bold)] for markdown strings with **bold** markers."""
        clean = normalize_markdown_inline(text)
        parts = re.split(r"(\*\*.*?\*\*)", clean)
        segments: list[tuple[str, bool]] = []
        for part in parts:
            if not part:
                continue
            if part.startswith("**") and part.endswith("**") and len(part) > 4:
                segments.append((part[2:-2], True))
            else:
                segments.append((part, False))
        return segments

    def add_markdown_paragraph(doc: Document, text: str, style: str | None = None):
        paragraph = doc.add_paragraph(style=style) if style else doc.add_paragraph()
        for segment, is_bold in markdown_segments(text):
            run = paragraph.add_run(segment)
            run.bold = is_bold
        return paragraph

    def add_plain_paragraph(doc: Document, text: str, style: str | None = None) -> None:
        paragraph = doc.add_paragraph(style=style) if style else doc.add_paragraph()
        paragraph.add_run(normalize_markdown_inline(text))

    def enable_toc_update_on_open(doc: Document) -> None:
        settings = doc.settings.element
        node = settings.find(qn("w:updateFields"))
        if node is None:
            node = OxmlElement("w:updateFields")
            settings.append(node)
        node.set(qn("w:val"), "true")

    def update_header_fields(doc: Document, tool_version: str) -> None:
        today = date.today().isoformat()

        def local_name(tag: str) -> str:
            return tag.split("}")[-1]

        for section in doc.sections:
            headers = [section.header, section.first_page_header, section.even_page_header]
            for header in headers:
                header_el = header._element

                # Update "Document Number" value (default template value: "PA")
                for table in header_el.iter():
                    if local_name(table.tag) != "tbl":
                        continue

                    rows = [node for node in table if local_name(node.tag) == "tr"]
                    for row_idx, row in enumerate(rows):
                        cells = [node for node in row if local_name(node.tag) == "tc"]
                        for col_idx, cell in enumerate(cells):
                            cell_text = "".join(node.text or "" for node in cell.iter() if local_name(node.tag) == "t")
                            if "Document Number" not in cell_text:
                                continue

                            if row_idx + 1 >= len(rows):
                                continue

                            next_row_cells = [node for node in rows[row_idx + 1] if local_name(node.tag) == "tc"]
                            if col_idx >= len(next_row_cells):
                                continue

                            for node in next_row_cells[col_idx].iter():
                                if local_name(node.tag) == "t" and (node.text or "").strip() == "PA":
                                    node.text = tool_version
                                    break
                            break

                for sdt in header_el.iter():
                    if local_name(sdt.tag) != "sdt":
                        continue

                    alias_name = None
                    for node in sdt.iter():
                        if local_name(node.tag) == "alias":
                            alias_name = node.get(qn("w:val"))
                            break

                    if alias_name not in {"Date", "Revision"}:
                        continue

                    alias_value = today if alias_name == "Date" else tool_version

                    for node in sdt.iter():
                        if local_name(node.tag) == "t":
                            node.text = alias_value
                            break

    def find_template_anchor(doc: Document):
        for paragraph in doc.paragraphs:
            if paragraph.text.strip().lower() == "(autogenerated-content-start)":
                return paragraph
        return None

    def remove_from_anchor_to_end(doc: Document, anchor_paragraph) -> None:
        body = doc._element.body
        start = body.index(anchor_paragraph._p)
        for idx in range(len(body) - 1, start - 1, -1):
            el = body[idx]
            # Keep section properties, otherwise doc.sections becomes empty and add_table() crashes
            if str(el.tag).endswith("}sectPr"):
                continue
            body.remove(el)

    def set_list_level_from_style(paragraph, level: int) -> None:
        """
        Force the list nesting level (w:ilvl) while keeping the list definition (w:numId)
        coming from the paragraph style (template-linked multilevel list).
        """
        level = max(0, min(int(level), 8))

        p = paragraph._p
        pPr = p.get_or_add_pPr()

        numPr = pPr.find(qn("w:numPr"))
        if numPr is None:
            numPr = OxmlElement("w:numPr")
            pPr.append(numPr)

        ilvl = numPr.find(qn("w:ilvl"))
        if ilvl is None:
            ilvl = OxmlElement("w:ilvl")
            numPr.append(ilvl)
        ilvl.set(qn("w:val"), str(level))

        # Ensure numId is present so the paragraph is attached to the same multilevel list as the style
        numId = numPr.find(qn("w:numId"))
        if numId is None:
            style_el = paragraph.style._element
            pPr_style = style_el.find(qn("w:pPr"))
            if pPr_style is None:
                return
            numPr_style = pPr_style.find(qn("w:numPr"))
            if numPr_style is None:
                return
            numId_style = numPr_style.find(qn("w:numId"))
            if numId_style is None:
                return

            numId = OxmlElement("w:numId")
            numId.set(qn("w:val"), numId_style.get(qn("w:val")))
            numPr.append(numId)


    # ------------------------------ Word logic ------------------------------ #
    text = md_file.read_text(encoding="utf-8")
    lines = text.splitlines()

    template_path = first_existing_path(DOCX_TEMPLATE_CANDIDATES)
    doc = Document(str(template_path)) if template_path else Document()

    anchor = find_template_anchor(doc)
    if anchor is not None:
        remove_from_anchor_to_end(doc, anchor)

    seen_heading1 = False  # Track if we've already written the first Heading 1 (## ...)

    i = 0
    while i < len(lines):
        line = lines[i]
        s = line.rstrip()

        if not s:
            # Skip blank lines around headings and around horizontal rules (---)
            def is_heading(line_: str) -> bool:
                t = (line_ or "").strip()
                return t.startswith("# ") or t.startswith("## ") or t.startswith("### ") or t.startswith("#### ")

            def is_hr(line_: str) -> bool:
                return (line_ or "").strip() == "---"

            # Previous non-empty line
            j = i - 1
            while j >= 0 and not lines[j].strip():
                j -= 1
            prev = lines[j] if j >= 0 else ""

            # Next non-empty line
            k = i + 1
            while k < len(lines) and not lines[k].strip():
                k += 1
            nxt = lines[k] if k < len(lines) else ""

            # If blank line is before/after a heading or an '---', drop it
            if is_heading(prev) or is_heading(nxt) or is_hr(prev) or is_hr(nxt):
                i += 1
                continue

            i += 1
            continue

        if s == "---":
            i += 1
            # Also skip any blank lines immediately after the horizontal rule
            while i < len(lines) and not lines[i].strip():
                i += 1
            continue

        if is_html_comment_line(s):
            i += 1
            continue

        image_info = parse_markdown_image_line(s.strip())
        if image_info:
            _, image_src = image_info
            image_path = resolve_markdown_image_path(md_file, image_src)
            if image_path:
                try:
                    doc.add_picture(str(image_path), width=DocxCm(16))
                except Exception:
                    pass
            i += 1
            continue

        t = s.strip()
        if t.startswith("| ") and t.endswith(" |") and i + 1 < len(lines) and is_table_separator(lines[i + 1]):
            rows: list[list[str]] = [parse_table_row(t)]
            i += 2
            while i < len(lines):
                candidate = lines[i].strip()
                if not (candidate.startswith("| ") and candidate.endswith(" |")):
                    break
                rows.append(parse_table_row(candidate))
                i += 1

            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = "Table Grid"
                for row_idx, row in enumerate(rows):
                    for col_idx, cell in enumerate(row):
                        paragraph = table.cell(row_idx, col_idx).paragraphs[0]
                        paragraph.text = ""  # limpia el párrafo de la celda
                        font_pt = WORD_TABLE_HEADER_PT if row_idx == 0 else WORD_TABLE_BODY_PT
                        for seg, is_bold in markdown_segments(cell):
                            run = paragraph.add_run(seg)
                            run.bold = is_bold or row_idx == 0
                            run.font.size = Pt(font_pt)

            continue

        if s.startswith("# "):
            # Markdown H1 is the document title (not part of heading numbering hierarchy).
            add_plain_paragraph(doc, clean_heading_text(s[2:].strip()), style="Title")

        elif s.startswith("## "):
            # Each Heading 1 starts on a new page, except the first one (usually right after the Title)
            if seen_heading1:
                doc.add_page_break()
            add_plain_paragraph(doc, clean_heading_text(s[3:].strip()), style="Heading 1")
            seen_heading1 = True

        elif s.startswith("### "):
            add_plain_paragraph(doc, clean_heading_text(s[4:].strip()), style="Heading 2")

        elif s.startswith("#### "):
            add_plain_paragraph(doc, clean_heading_text(s[5:].strip()), style="Heading 3")

        # Bullets (support sub-bullets by indentation)
        elif re.match(r"^\s*-\s+", s):
            m = re.match(r"^(\s*)-\s+(.*)$", s)
            indent = len(m.group(1).replace("\t", "    "))
            item_text = m.group(2).strip()

            # 2 spaces per level (as in your example). Use //4 if your markdown uses 4 spaces per level.
            level = 0 if indent < 2 else min(8, indent // 2)

            paragraph = doc.add_paragraph(style=WORD_LIST_STYLE)
            for segment, is_bold in markdown_segments(item_text):
                run = paragraph.add_run(segment)
                run.bold = is_bold

            set_list_level_from_style(paragraph, level)



        elif re.match(r"^[A-Za-z](?:\.[A-Za-z])*\)\s+", s):
            # Examples: A) , B) , a) , A.B) , a.b)
            add_markdown_paragraph(doc, re.sub(r"^[A-Za-z](?:\.[A-Za-z])*\)\s+", "", s), style="List Number")

        elif re.match(r"^\d+\.\s+", s):
            # Keep Markdown numbering to avoid Word continuing previous lists
            add_markdown_paragraph(doc, s.strip(), style="List Paragraph")

        else:
            add_markdown_paragraph(doc, s)

        i += 1

    enable_toc_update_on_open(doc)
    update_header_fields(doc, version)
    doc.save(docx_file)


# Markdown to PowerPoint Converter
def build_pptx_summary(md_file: Path, pptx_file: Path, version: str) -> None:
    print(f"\tGenerating PowerPoint from Markdown...")
    # ------------------- PPT-only helpers (subfunctions) ------------------- #
    def _strip_bold(text: str) -> str:
        return re.sub(r"\*\*(.*?)\*\*", r"\1", text)

    def _norm(text: str) -> str:
        return normalize_markdown_inline(_strip_bold(text))

    def _indent_to_level(indent_spaces: int) -> int:
        # 0-1 => 0, 2-7 => 1, 8-11 => 2, ...
        if indent_spaces < 2:
            return 0
        lvl = indent_spaces // 4
        if lvl == 0:
            lvl = 1
        return min(lvl, 4)

    def _parse_sections_for_ppt(md_path: Path) -> list[dict]:
        """
        Parse markdown into sections (##) and blocks.
        - ### becomes block type 'h3' (we'll use it as slide subtitle)
        - #### becomes block type 'h4'
        - bullets preserve indentation
        - supports numeric ordered (1.) and alpha ordered (A), A.B))
        """
        lines = md_path.read_text(encoding="utf-8").splitlines()
        sections: list[dict] = []
        current = {"title": "Overview", "blocks": []}

        bullet_pat = re.compile(r"^(\s*)-\s+(.*)$")
        num_pat = re.compile(r"^(\s*)(\d+)\.\s+(.*)$")
        alpha_pat = re.compile(r"^(\s*)([A-Za-z](?:\.[A-Za-z])*)\)\s+(.*)$")

        i = 0
        while i < len(lines):
            raw = lines[i].rstrip()
            line = raw.strip()

            if not line or line == "---":
                i += 1
                continue

            if is_html_comment_line(line):
                i += 1
                continue

            if line.startswith("# "):
                i += 1
                continue

            if line.startswith("## "):
                if current["blocks"]:
                    sections.append(current)
                current = {"title": clean_heading_text(line[3:].strip()), "blocks": []}
                i += 1
                continue

            if line.startswith("### "):
                current["blocks"].append({"type": "h3", "text": clean_heading_text(line[4:].strip())})
                i += 1
                continue

            if line.startswith("#### "):
                current["blocks"].append({"type": "h4", "text": clean_heading_text(line[5:].strip())})
                i += 1
                continue

            if line.startswith("| ") and line.endswith(" |") and i + 1 < len(lines) and is_table_separator(lines[i + 1]):
                rows = [parse_table_row(line)]
                i += 2
                while i < len(lines):
                    cand_raw = lines[i].rstrip()
                    cand = cand_raw.strip()
                    if not (cand.startswith("| ") and cand.endswith(" |")):
                        break
                    rows.append(parse_table_row(cand))
                    i += 1
                current["blocks"].append({"type": "table", "rows": rows})
                continue

            image_info = parse_markdown_image_line(line)
            if image_info:
                alt_text, image_src = image_info
                image_path = resolve_markdown_image_path(md_path, image_src)
                if image_path:
                    current["blocks"].append({
                        "type": "image",
                        "alt": _norm(alt_text) or "Screenshot",
                        "path": str(image_path),
                    })
                i += 1
                continue

            mb = bullet_pat.match(raw)
            mn = num_pat.match(raw)
            ma = alpha_pat.match(raw)

            if mb or mn or ma:
                items: list[dict] = []
                mode = "bullet" if mb else ("num" if mn else "alpha")

                while i < len(lines):
                    raw_c = lines[i].rstrip("\n")
                    c = raw_c.strip()

                    if not c or c == "---":
                        break
                    if c.startswith("## ") or c.startswith("### ") or c.startswith("#### "):
                        break
                    if c.startswith("| "):
                        break

                    mb2 = bullet_pat.match(raw_c)
                    mn2 = num_pat.match(raw_c)
                    ma2 = alpha_pat.match(raw_c)

                    if mode == "bullet" and mb2:
                        indent_spaces = len(mb2.group(1).replace("\t", "    "))
                        items.append({
                            "text": _norm(mb2.group(2).strip()),
                            "level": _indent_to_level(indent_spaces),
                            "marker": None,
                        })
                        i += 1
                        continue

                    if mode == "num" and mn2:
                        indent_spaces = len(mn2.group(1).replace("\t", "    "))
                        marker = f"{mn2.group(2)}."
                        items.append({
                            "text": _norm(mn2.group(3).strip()),
                            "level": _indent_to_level(indent_spaces),
                            "marker": marker,
                        })
                        i += 1
                        continue

                    if mode == "alpha" and ma2:
                        indent_spaces = len(ma2.group(1).replace("\t", "    "))
                        marker = f"{ma2.group(2)})"
                        items.append({
                            "text": _norm(ma2.group(3).strip()),
                            "level": _indent_to_level(indent_spaces),
                            "marker": marker,
                        })
                        i += 1
                        continue

                    break

                current["blocks"].append({
                    "type": "list",
                    "ordered": mode in ("num", "alpha"),
                    "items": items,
                })
                continue

            current["blocks"].append({"type": "paragraph", "text": _norm(line)})
            i += 1

        if current["blocks"]:
            sections.append(current)

        return sections

    def _block_units(block: dict) -> int:
        # Rough estimator of "lines" to fit content in a slide.
        # Calibrated to a 20-lines max template.
        if block["type"] == "h4":
            return 2
        if block["type"] == "paragraph":
            return max(1, ceil(len(block["text"]) / PPT_CUT_SLIDE_WEIGHT))
        if block["type"] == "list":
            return 1 + sum(max(1, ceil(len(it["text"]) / PPT_CUT_SLIDE_WEIGHT)) for it in block["items"])
        if block["type"] == "table":
            return 6 + len(block["rows"])
        if block["type"] == "image":
            return 14
        return 2

    def _style_table(table) -> None:
        header_fill = RGBColor(31, 78, 121)
        band_fill = RGBColor(242, 246, 252)
        text_dark = RGBColor(40, 40, 40)

        for col in range(len(table.columns)):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(PPT_FONT_SIZE_TABLES_HEADER)

        for row in range(1, len(table.rows)):
            for col in range(len(table.columns)):
                cell = table.cell(row, col)
                if row % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = band_fill
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = text_dark
                p.font.size = Pt(PPT_FONT_SIZE_TABLES_BODY)

    def _set_table_wrap(table) -> None:
        for row in table.rows:
            for cell in row.cells:
                tf = cell.text_frame
                tf.word_wrap = True
                tf.margin_left = Pt(2)
                tf.margin_right = Pt(2)
                tf.margin_top = Pt(1)
                tf.margin_bottom = Pt(1)

    def _autosize_table(table, total_w, max_h, body_font_pt=9, header_font_pt=10,
                       min_col_in=0.9, max_col_in=3.2) -> int:
        EMU_PER_INCH = 914400
        EMU_PER_PT = 12700

        nrows = len(table.rows)
        ncols = len(table.columns)

        max_lens = [1] * ncols
        for r in range(nrows):
            for c in range(ncols):
                txt = (table.cell(r, c).text or "").strip()
                max_lens[c] = max(max_lens[c], len(txt))

        weights = [max(1.0, ml ** 0.6) for ml in max_lens]
        wsum = sum(weights) or 1.0

        min_w = int(min_col_in * EMU_PER_INCH)
        max_w = int(max_col_in * EMU_PER_INCH)

        widths = [int(total_w * (w / wsum)) for w in weights]
        widths = [max(min_w, min(max_w, cw)) for cw in widths]

        diff = int(total_w - sum(widths))
        widths[-1] = max(min_w, min(max_w, widths[-1] + diff))
        diff2 = int(total_w - sum(widths))
        widths[-1] += diff2

        for c, cw in enumerate(widths):
            table.columns[c].width = cw

        def _cap_chars(col_w_emu: int, font_pt: float) -> int:
            col_pt = col_w_emu / EMU_PER_PT
            return max(6, int(col_pt / (0.55 * font_pt)))

        def _est_lines(txt: str, cap: int) -> int:
            if not txt:
                return 1
            lines = 0
            for part in txt.splitlines():
                lines += max(1, ceil(len(part) / cap))
            return lines

        total_h = 0
        for r in range(nrows):
            font_pt = header_font_pt if r == 0 else body_font_pt
            max_lines = 1
            for c in range(ncols):
                cap = _cap_chars(widths[c], font_pt)
                txt = (table.cell(r, c).text or "").strip()
                max_lines = max(max_lines, _est_lines(txt, cap))

            # row_h_pt = max_lines * (font_pt * 1.2) + 6
            row_h_pt = max_lines * (font_pt * 1.05) + 2

            table.rows[r].height = Pt(row_h_pt)
            total_h += table.rows[r].height

        return total_h

    def _try_add_small_table_below_text(
        prs: Presentation,
        slide_title: str,
        subtitle: str | None,
        continuation: bool,
        buffer: list[dict],
        used: int,
        max_units: int,
        rows: list[list[str]],
        combine_max_data_rows: int,
        gap,
        min_text_h,
        min_table_h,
    ) -> tuple[bool, list[dict], int, bool]:
        """
        Try to place a small table below the current buffered text in the same slide.
        This is heuristic: it uses 'used/max_units' to estimate how much vertical space the text needs.
        Returns: (handled, new_buffer, new_used, new_continuation)
        """
        if not buffer or not rows or len(rows) < 2:
            return False, buffer, used, continuation

        header = rows[:1]
        data = rows[1:]
        if len(data) > combine_max_data_rows:
            return False, buffer, used, continuation

        # Estimate available geometry from the layout (no slide created yet).
        ly = prs.slide_layouts[1]
        ph = ly.placeholders[1]
        x0, y0, w0, h0 = ph.left, ph.top, ph.width, ph.height

        # Estimate text height proportionally to 'used/max_units' and clamp it.
        text_h = int(h0 * min(used, max_units) / max_units)
        if text_h < int(min_text_h):
            text_h = int(min_text_h)

        max_text_h = int(h0 - min_table_h - gap)
        if text_h > max_text_h:
            return False, buffer, used, continuation

        tbl_y = y0 + text_h + gap
        tbl_h = (y0 + h0) - tbl_y
        if tbl_h < int(min_table_h):
            return False, buffer, used, continuation

        # Create the slide only if we can actually fit the table.
        slide_local = prs.slides.add_slide(prs.slide_layouts[1])

        # Title + subtitle (same style as normal slides)
        title_shape = slide_local.shapes.title
        title_shape.text = slide_title
        if subtitle:
            p = title_shape.text_frame.add_paragraph()
            p.text = subtitle if not continuation else f"{subtitle} (cont.)"
            p.level = 0
            p.font.bold = False
            p.font.size = Pt(20)
            p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

        content = slide_local.shapes.placeholders[1]
        x, y, w, h = content.left, content.top, content.width, content.height

        # Render buffered text in the template placeholder (keeps template styles).
        # IMPORTANT: do NOT resize the placeholder height, it can trigger per-character wrapping in some templates.
        _render_content(content.text_frame, buffer)

        # Add table below (table is added after the text so it sits on top if shapes overlap).
        chunk = header + data
        tbl_shape = slide_local.shapes.add_table(len(chunk), len(chunk[0]), x, tbl_y, w, tbl_h)

        table = tbl_shape.table

        for r, row in enumerate(chunk):
            for c, value in enumerate(row):
                table.cell(r, c).text = _norm(value)

        _style_table(table)
        _set_table_wrap(table)

        # Autosize rows/cols and shrink table height so it doesn't occupy all the area
        total_h = _autosize_table(table, total_w=w, max_h=tbl_h, body_font_pt=9, header_font_pt=10)
        tbl_shape.height = min(total_h, tbl_h)
        tbl_shape.top = tbl_y

        # We consumed the buffered text + this table
        return True, [], 0, True


    def _add_table_slide(prs: Presentation, slide_title: str, subtitle: str | None, rows: list[list[str]], continuation: bool = False) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_shape = slide.shapes.title
        title_shape.text = slide_title
        if subtitle:
            p = title_shape.text_frame.add_paragraph()
            p.text = subtitle if not continuation else f"{subtitle} (cont.)"
            p.level = 0
            p.font.bold = False
            p.font.size = Pt(PPT_FONT_SIZE_H3)
            p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

        content = slide.shapes.placeholders[1]
        x, y, w, h = content.left, content.top, content.width, content.height

        # Important: don't leave it empty, otherwise PowerPoint shows "Click to add text"
        tf = content.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = " "
        p0.font.size = Pt(1)
        p0.font.color.rgb = RGBColor(255, 255, 255)

        tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
        table = tbl_shape.table

        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                table.cell(r, c).text = _norm(value)

        _style_table(table)
        _set_table_wrap(table)

        total_h = _autosize_table(table, total_w=w, max_h=h, body_font_pt=9, header_font_pt=10)
        tbl_shape.height = min(total_h, h)
        tbl_shape.top = y

    def _content_placeholders_in_order(slide) -> list:
        placeholders = []
        for ph in slide.shapes.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                continue
            placeholders.append(ph)
        placeholders.sort(key=lambda shape: (shape.top, shape.left))
        return placeholders

    def _add_image_slide(
        prs: Presentation,
        slide_title: str,
        subtitle: str | None,
        image_path: str,
        caption: str | None,
        continuation: bool = False,
        second_image_path: str | None = None,
        text_blocks: list[dict] | None = None,
    ) -> None:
        has_second_image = bool(second_image_path)
        has_text_blocks = bool(text_blocks)
        # If the section contains only a single image, use layout 1.
        if has_second_image:
            layout_idx = 3
        elif has_text_blocks:
            layout_idx = 2
        else:
            layout_idx = 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        title_shape = slide.shapes.title
        title_shape.text = slide_title
        if subtitle:
            p = title_shape.text_frame.add_paragraph()
            p.text = subtitle if not continuation else f"{subtitle} (cont.)"
            p.level = 0
            p.font.bold = False
            p.font.size = Pt(PPT_FONT_SIZE_H3)
            p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

        content_placeholders = _content_placeholders_in_order(slide)
        if has_second_image:
            needed_placeholders = 3
        elif has_text_blocks:
            needed_placeholders = 2
        else:
            needed_placeholders = 1
        if len(content_placeholders) < needed_placeholders:
            raise RuntimeError(
                f"Layout {layout_idx} must provide at least {needed_placeholders} content placeholders; found {len(content_placeholders)}"
            )

        first_image_placeholder = content_placeholders[0]
        slide.shapes.add_picture(
            image_path,
            first_image_placeholder.left,
            first_image_placeholder.top,
            width=first_image_placeholder.width,
            height=first_image_placeholder.height,
        )

        if has_second_image:
            second_image_placeholder = content_placeholders[1]
            slide.shapes.add_picture(
                second_image_path,
                second_image_placeholder.left,
                second_image_placeholder.top,
                width=second_image_placeholder.width,
                height=second_image_placeholder.height,
            )
            text_placeholder = content_placeholders[2]
        elif has_text_blocks:
            text_placeholder = content_placeholders[1]
        else:
            text_placeholder = None

        if text_placeholder is not None:
            tf = text_placeholder.text_frame
            if text_blocks:
                _render_content(tf, text_blocks)
            else:
                tf.clear()
                p0 = tf.paragraphs[0]
                p0.text = caption if caption else " "
                if caption:
                    p0.level = 0
                    p0.font.size = Pt(10)
                    p0.font.bold = True
                else:
                    p0.font.size = Pt(1)
                    p0.font.color.rgb = RGBColor(255, 255, 255)
        elif caption:
            # No text placeholder available (image-only layout), so append caption to title.
            p = title_shape.text_frame.add_paragraph()
            p.text = caption
            p.level = 0
            p.font.bold = False
            p.font.size = Pt(10)

    def _render_content(tf, blocks: list[dict]) -> None:
        from pptx.oxml.ns import qn as pptx_qn
        from pptx.oxml.xmlchemy import OxmlElement as PptxOxmlElement

        def _set_bullets(p, enabled: bool) -> None:
            pPr = p._p.get_or_add_pPr()
            for tag in ("a:buAutoNum", "a:buChar", "a:buBlip", "a:buFont", "a:buNone"):
                el = pPr.find(pptx_qn(tag))
                if el is not None:
                    pPr.remove(el)
            if enabled:
                bu = PptxOxmlElement("a:buChar")
                bu.set("char", "•")
                pPr.insert(0, bu)
            else:
                pPr.insert(0, PptxOxmlElement("a:buNone"))

        tf.clear()
        first = True

        for block in blocks:
            if block["type"] == "h4":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = block["text"]
                p.level = 0
                _set_bullets(p, enabled=False)
                p.font.bold = True
                p.font.size = Pt(PPT_FONT_SIZE_H4)
                continue

            if block["type"] == "paragraph":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = block["text"]
                p.level = 0
                _set_bullets(p, enabled=False)
                p.font.size = Pt(PPT_FONT_SIZE_PARAGRAPH)
                continue

            if block["type"] == "list":
                for idx, it in enumerate(block["items"], start=1):
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False

                    nesting = int(it.get("level", 0))
                    p.level = 1 + nesting
                    p.font.size = Pt(PPT_FONT_SIZE_LISTS)

                    if block["ordered"]:
                        marker = it.get("marker") or f"{idx}."
                        p.text = f"{marker} {it['text']}"
                        _set_bullets(p, enabled=False)
                    else:
                        p.text = it["text"]
                        _set_bullets(p, enabled=True)

    # ------------------------------ PPT logic ------------------------------ #
    sections = _parse_sections_for_ppt(md_file)

    template_path = first_existing_path(PPTX_TEMPLATE_CANDIDATES)
    prs = Presentation(str(template_path)) if template_path else Presentation()

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Technical User Guide"
    slide.placeholders[1].text = f"SSB Retuning Automations v{version}"

    for section in sections:
        slide_title = section["title"]
        blocks = section["blocks"]

        groups: list[dict] = []
        current_subtitle: str | None = None
        current_blocks: list[dict] = []

        for b in blocks:
            if b["type"] == "h3":
                if current_subtitle is not None or current_blocks:
                    groups.append({"subtitle": current_subtitle, "blocks": current_blocks})
                current_subtitle = b["text"]
                current_blocks = []
            else:
                current_blocks.append(b)

        if current_subtitle is not None or current_blocks:
            groups.append({"subtitle": current_subtitle, "blocks": current_blocks})

        for g in groups:
            subtitle = g["subtitle"]
            gblocks = g["blocks"]

            buffer: list[dict] = []
            used = 0
            continuation = False

            def flush() -> None:
                nonlocal buffer, used, continuation
                slide_local = prs.slides.add_slide(prs.slide_layouts[1])

                title_shape = slide_local.shapes.title
                title_shape.text = slide_title
                if subtitle:
                    p = title_shape.text_frame.add_paragraph()
                    p.text = subtitle if not continuation else f"{subtitle} (cont.)"
                    p.level = 0
                    p.font.bold = False
                    p.font.size = Pt(20)
                    p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

                _render_content(slide_local.shapes.placeholders[1].text_frame, buffer)

                buffer = []
                used = 0
                continuation = True

            i = 0
            while i < len(gblocks):
                b = gblocks[i]

                if b["type"] == "table":
                    rows = b["rows"] or []
                    if not rows:
                        i += 1
                        continue

                    # Try to place a small table below the current buffered text (same slide)
                    handled, buffer, used, continuation = _try_add_small_table_below_text(
                        prs=prs,
                        slide_title=slide_title,
                        subtitle=subtitle,
                        continuation=continuation,
                        buffer=buffer,
                        used=used,
                        max_units=PPT_MAX_CONTENT_LINES,
                        rows=rows,
                        combine_max_data_rows=PPT_COMBINE_MAX_DATA_ROWS,
                        gap=PPT_GAP,
                        min_text_h=PPT_MIN_TEXT_H,
                        min_table_h=PPT_MIN_TABLE_H,
                    )
                    if handled:
                        i += 1
                        continue

                    # Normal behavior: keep order stable (text slide first, then table slides)
                    if buffer:
                        flush()

                    header = rows[:1]
                    data = rows[1:] if len(rows) > 1 else []

                    if not data:
                        _add_table_slide(prs, slide_title, subtitle, rows, continuation=continuation)
                        continuation = True
                    else:
                        for start in range(0, len(data), PPT_MAX_TABLE_DATA_ROWS):
                            chunk = header + data[start:start + PPT_MAX_TABLE_DATA_ROWS]
                            _add_table_slide(
                                prs,
                                slide_title,
                                subtitle,
                                chunk,
                                continuation=(continuation or start > 0),
                            )
                        continuation = True

                    i += 1
                    continue

                if b["type"] == "image":
                    if buffer:
                        flush()

                    second_image_path = None
                    text_after_images: list[dict] = []

                    if i + 1 < len(gblocks) and gblocks[i + 1]["type"] == "image":
                        second_image_path = gblocks[i + 1]["path"]
                        j = i + 2
                    else:
                        j = i + 1

                    while j < len(gblocks) and gblocks[j]["type"] in {"paragraph", "list", "h4"}:
                        text_after_images.append(gblocks[j])
                        j += 1

                    _add_image_slide(
                        prs,
                        slide_title,
                        subtitle,
                        b["path"],
                        b.get("alt"),
                        continuation=continuation,
                        second_image_path=second_image_path,
                        text_blocks=text_after_images,
                    )
                    continuation = True
                    i = j
                    continue

                # Special case: split long lists so they respect the per-slide line budget
                if b["type"] == "list":
                    u = _block_units(b)
                    remaining = PPT_MAX_CONTENT_LINES - used

                    if u > PPT_MAX_CONTENT_LINES or (buffer and used + u > PPT_MAX_CONTENT_LINES):
                        # Need at least "1 unit" overhead + 1 item to be worth placing on this slide
                        fit_items = []
                        fit_used = 1  # list overhead (same idea as _block_units)
                        for it in b["items"]:
                            item_u = max(1, ceil(len(it["text"]) / PPT_CUT_SLIDE_WEIGHT))
                            if fit_used + item_u > remaining:
                                break
                            fit_items.append(it)
                            fit_used += item_u

                        if fit_items:
                            # Put the first part of the list on the current slide
                            b1 = dict(b)
                            b1["items"] = fit_items
                            buffer.append(b1)
                            used += fit_used

                            # Flush now so the list continues on the next slide
                            flush()

                            # Replace current block with the remaining items and re-process it
                            rest = b["items"][len(fit_items):]
                            if rest:
                                gblocks[i] = {"type": "list", "ordered": b["ordered"], "items": rest}
                                continue

                            i += 1
                            continue

                        # If nothing fits (e.g., a single very long item), force one item per slide
                        if b["items"]:
                            b1 = dict(b)
                            b1["items"] = [b["items"][0]]
                            buffer.append(b1)
                            used = PPT_MAX_CONTENT_LINES
                            flush()

                            rest = b["items"][1:]
                            if rest:
                                gblocks[i] = {"type": "list", "ordered": b["ordered"], "items": rest}
                                continue

                            i += 1
                            continue

                        flush()
                        continue

                # Default behavior for non-list blocks (and lists that fit)
                u = _block_units(b)
                if buffer and used + u > PPT_MAX_CONTENT_LINES:
                    flush()

                buffer.append(b)
                used += u
                i += 1

            if buffer:
                flush()

    prs.save(pptx_file)



def try_export_pptx_pdf_windows(pptx_file: Path, pdf_file: Path) -> bool:
    """Best effort on Windows: export PPTX to PDF through PowerPoint automation."""
    if os.name != "nt":
        return False

    ps_script = f"""
$ErrorActionPreference = 'Stop'
$ppt = $null
try {{
    $pptPath = '{str(pptx_file)}'
    $pdfPath = '{str(pdf_file)}'

    $ppt = New-Object -ComObject PowerPoint.Application
    $presentation = $ppt.Presentations.Open($pptPath, $false, $false, $false)
    # 32 = ppSaveAsPDF
    $presentation.SaveAs($pdfPath, 32)
    $presentation.Close()
}} catch {{
    Write-Error $_.Exception.Message
    exit 1
}} finally {{
    if ($ppt -ne $null) {{ $ppt.Quit() }}
}}
"""

    result = subprocess.run(
        ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps_script],
        cwd=str(ROOT),
        text=True,
        encoding="utf-8",
        errors="replace",
        capture_output=True,
    )
    return result.returncode == 0 and pdf_file.exists()


def build_pdf_from_pptx(pptx_file: Path, pdf_file: Path) -> bool:
    """Best-effort PPTX -> PDF conversion using PowerPoint (Windows) or LibreOffice."""
    print("	Converting PPTX to PDF...")
    if not pptx_file.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_file}")

    output_dir = pdf_file.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    if try_export_pptx_pdf_windows(pptx_file, pdf_file):
        return True

    commands: list[list[str]] = []
    soffice = shutil.which("soffice")
    libreoffice = shutil.which("libreoffice")
    if soffice:
        commands.append([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_file)])
    if libreoffice:
        commands.append([libreoffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_file)])

    generated_tmp_pdf = output_dir / f"{pptx_file.stem}.pdf"
    for cmd in commands:
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, check=False)
            if result.returncode == 0 and generated_tmp_pdf.exists():
                generated_tmp_pdf.replace(pdf_file)
                return True
        except OSError:
            continue

    print("	[WARN] Unable to convert PPTX to PDF (PowerPoint/LibreOffice not available).")
    return False


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Update user guide artifacts.")
    parser.add_argument(
        "--formats",
        nargs="+",
        choices=["docx", "docx.pdf", "pptx", "pptx.pdf"],
        default=["docx", "docx.pdf", "pptx", "pptx.pdf"],
        help="Guide formats to generate/update. Default: all.",
    )
    return parser.parse_args()


def update_user_guides(formats: set[str]) -> dict[str, Path]:
    tool_version = get_tool_version()
    print(f"🔍 TOOL_VERSION detected: v{tool_version}")

    print(f"▶️ Updating README.md User Guides links to new version: {tool_version}...")
    update_readme_links(tool_version)

    print(f"▶️ Updating User Guides to new version: {tool_version}...")
    paths = align_help_guides_to_version(tool_version)

    if "docx" in formats:
        build_docx_from_markdown(paths["md"], paths["docx"], tool_version)

    if "pptx" in formats:
        build_pptx_summary(paths["md"], paths["pptx"], tool_version)

    if "docx.pdf" in formats:
        if paths["docx"].exists() and try_update_docx_fields_and_export_pdf(paths["docx"], paths["docx_pdf"]):
            pass
        else:
            build_pdf_from_markdown(paths["md"], paths["docx_pdf"])

    pptx_pdf_generated = False
    if "pptx.pdf" in formats:
        if not paths["pptx"].exists():
            # PDF conversion requires a PPTX source; generate it when missing.
            build_pptx_summary(paths["md"], paths["pptx"], tool_version)
        before_exists = paths["pptx_pdf"].exists()
        before_mtime = paths["pptx_pdf"].stat().st_mtime if before_exists else None
        pptx_pdf_generated = build_pdf_from_pptx(paths["pptx"], paths["pptx_pdf"])
        after_exists = paths["pptx_pdf"].exists()
        after_mtime = paths["pptx_pdf"].stat().st_mtime if after_exists else None
        if (not pptx_pdf_generated) and (not after_exists):
            raise RuntimeError(
                f"Failed to generate {paths['pptx_pdf']}. "
                "Install Microsoft PowerPoint (Windows) or LibreOffice/soffice in PATH."
            )
        if after_exists and before_exists and before_mtime == after_mtime and not pptx_pdf_generated:
            print("	[WARN] PPTX PDF conversion did not refresh the existing file; keeping previous artifact.")

    cleanup_old_versioned_guides(paths)

    print(f"Using markdown: {paths['md']}")
    if "docx" in formats:
        print(f"Generated: {paths['docx']}")
    if "pptx" in formats:
        print(f"Generated: {paths['pptx']}")
    if "docx.pdf" in formats:
        print(f"Generated: {paths['docx_pdf']}")
    if "pptx.pdf" in formats:
        if pptx_pdf_generated:
            print(f"Generated: {paths['pptx_pdf']}")
        elif paths["pptx_pdf"].exists():
            print(f"Generated (existing): {paths['pptx_pdf']}")
        else:
            print(f"Skipped (not generated): {paths['pptx_pdf']}")
    print("User Guides updated.")
    print("README Technical Guide Links updated.")
    return paths


if __name__ == "__main__":
    args = parse_args()
    update_user_guides(set(args.formats))
