from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
md_path = ROOT / "help" / "SSB_RetuningAutomations_Technical_User_Guide_Detailed.md"
docx_path = ROOT / "help" / "SSB_RetuningAutomations_Technical_User_Guide_Detailed.docx"
pptx_path = ROOT / "help" / "SSB_RetuningAutomations_Technical_User_Guide_Resumen.pptx"

text = md_path.read_text(encoding="utf-8")
lines = text.splitlines()

# ---- DOCX ----
doc = Document()
for line in lines:
    s = line.rstrip()
    if not s:
        doc.add_paragraph("")
        continue
    if s.startswith("# "):
        doc.add_heading(s[2:].strip(), level=1)
    elif s.startswith("## "):
        doc.add_heading(s[3:].strip(), level=2)
    elif s.startswith("### "):
        doc.add_heading(s[4:].strip(), level=3)
    elif s.startswith("- "):
        doc.add_paragraph(s[2:].strip(), style="List Bullet")
    elif s.startswith("| ") and s.endswith(" |"):
        # keep markdown tables readable as plain text
        doc.add_paragraph(s)
    else:
        doc.add_paragraph(s)

doc.save(docx_path)

# ---- PPTX (resumen) ----
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "SSB_RetuningAutomations"
slide.placeholders[1].text = "Technical User Guide (Resumen Ejecutivo)"

slides_content = [
    ("Objetivo de la herramienta", [
        "Automatizar auditorías y comparaciones pre/post en retuning SSB.",
        "Ejecución GUI o CLI con salida versionada y trazable.",
        "Módulos: Update frequencies, Configuration Audit, Consistency Check, Bulk, Clean-Up.",
    ]),
    ("Módulo 1: Configuration Audit", [
        "Entrada: logs/ZIP + frecuencias y listas permitidas.",
        "Salida: ConfigurationAudit.xlsx + ConfigurationAudit.pptx + comandos CA.",
        "Genera SummaryAudit con chequeos NR/LTE/ENDC/Perfiles/Cardinalidades.",
    ]),
    ("Summary Audit (qué evalúa)", [
        "MeContext: nodos totales y UNSYNCHRONIZED excluidos.",
        "NR/LTE: old/new SSB, inconsistencias y valores fuera de listas permitidas.",
        "ENDC/FreqPrioNR: coherencia old/new/N77B y mismatches de parámetros.",
        "Externals/TermPoints: trazabilidad de destinos SSB-Pre/SSB-Post/Unknown.",
    ]),
    ("Módulo 2: Consistency Check", [
        "Detecta relaciones nuevas, faltantes y discrepancias en relaciones comunes.",
        "Discrepancias de parámetros: comparación columna a columna (con exclusiones de ruido).",
        "Discrepancias de frecuencia: PRE tenía old/new y POST no quedó en SSB objetivo.",
        "Clasifica discrepancias: SSB-Post vs Unknown.",
    ]),
    ("Outputs de ConsistencyChecks", [
        "Summary, SummaryAuditComparisson y Summary_CellRelation.",
        "Hojas GU: relations, param_disc, freq_disc, freq_disc_unknown, missing, new.",
        "Hojas NR: relations, param_disc, freq_disc, freq_disc_unknown, missing, new.",
        "Correction_Cmd_CC con comandos para remediación.",
    ]),
    ("Módulo 3/4 y operación recomendada", [
        "Bulk: autodetección PRE/POST por fecha/hora y mercado.",
        "Final Clean-Up: base lista para políticas de cierre.",
        "Primero revisar Summary, luego detalle por discrepancia y comandos.",
    ]),
]

for title, bullets in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)

prs.save(pptx_path)
print(f"Generated: {docx_path}")
print(f"Generated: {pptx_path}")
