# -*- coding: utf-8 -*-

import os
from typing import Dict, List, Optional
import pandas as pd

from src.utils.utils_infrastructure import get_resource_path


def build_text_summary_structure(
    summary_audit_df: pd.DataFrame,
) -> Dict[str, List[Dict[str, object]]]:
    """
    Group SummaryAudit rows by top-level Category, keeping SubCategory/Metric/Value/ExtraInfo.

    Returns:
      {
        "NR Frequency Audit": [
            {"SubCategory": "NRCellDU", "Metric": "...", "Value": 4, "ExtraInfo": "..."},
            ...
        ],
        "NR Frequency Inconsistencies": [
            ...
        ],
        ...
      }
    """
    sections: Dict[str, List[Dict[str, object]]] = {}

    if summary_audit_df is None or summary_audit_df.empty:
        sections["Info"] = [
            {
                "SubCategory": "Info",
                "Metric": "No audit data available to build textual summary",
                "Value": "",
                "ExtraInfo": "",
            }
        ]
        return sections

    for _, row in summary_audit_df.iterrows():
        category = str(row.get("Category", "") or "Info")
        item = {
            "SubCategory": str(row.get("SubCategory", "") or ""),
            "Metric": str(row.get("Metric", "") or ""),
            "Value": row.get("Value", ""),
            "ExtraInfo": str(row.get("ExtraInfo", "") or ""),
        }
        sections.setdefault(category, []).append(item)

    return sections


def generate_ppt_summary(
    summary_audit_df: pd.DataFrame,
    excel_path: str,
    module_name: str = "",
) -> Optional[str]:
    """
    Generate a PPTX file next to the Excel with slides grouped by top-level Category.

    - First slide: global title.
    - Then, for each Category:

        • If Category name contains 'audit' (case-insensitive):
            - Single slide per Category.
            - Title = Category.
            - Body = one bullet per row: "Metric: Value" (no node list).

        • If Category name contains 'inconsist' (case-insensitive):
            - One or more slides per Category.
            - Title = Category.
            - For each row with Value > 0 and a non-empty node list:
                · Main bullet "Metric: Value".
                · Level-1 bullets with the node list parsed from ExtraInfo
                  (comma/semicolon separated), split into blocks of 50 items
                  per slide (no truncation).
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        print(f"{module_name} [INFO] python-pptx is not installed. Skipping PPT summary.")
        return None

    MAIN_BULLET_SIZE = Pt(14)
    SUB_BULLET_SIZE = Pt(10)

    sections = build_text_summary_structure(summary_audit_df)

    base, _ = os.path.splitext(excel_path)
    ppt_path = base + "_Summary.pptx"

    def _set_paragraph_font_size(paragraph, size: Pt) -> None:
        for run in paragraph.runs:
            run.font.size = size

    def _value_is_positive(v: object) -> bool:
        """Return True if the given value represents a numeric value > 0."""
        try:
            if isinstance(v, (int, float)):
                return v > 0
            s = str(v).strip()
            if not s or s.upper() == "N/A":
                return False
            return float(s) > 0
        except Exception:
            return False

    template_path = get_resource_path("ppt_templates/ConfigurationAuditTemplate.pptx")
    try:
        prs = Presentation(template_path)
        print(f"{module_name} Using PPT template: {template_path}")
    except Exception as e:
        print(f"{module_name} [WARN] Could not load PPT template, using default. ({e})")
        prs = Presentation()

    try:
        title_slide_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[2]
    except Exception:
        title_slide_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    # --- Title slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    title.text = "Configuration Audit Summary"
    if subtitle is not None:
        subtitle.text = os.path.basename(excel_path)

    # --- Category slides ---
    for category, items in sections.items():
        cat_lower = category.lower()
        is_audit = "audit" in cat_lower
        is_incons = "inconsist" in cat_lower  # covers 'Inconsistences' typo as well

        # ---------------------- INCONSISTENCIES: may need multiple slides ----------------------
        if is_incons:
            if not items:
                # No items at all: skip creating slides for this category
                continue

            for item in items:
                metric = item.get("Metric", "")
                value = item.get("Value", "")
                extra = item.get("ExtraInfo", "")

                # Skip rows with non-positive value
                if not _value_is_positive(value):
                    continue

                # Parse node/cell list from ExtraInfo
                cleaned_extra = str(extra).replace(";", ",") if extra else ""
                nodes = [t.strip() for t in cleaned_extra.split(",") if t.strip()]

                # Skip rows with empty node/cell list
                if not nodes:
                    continue

                main_text = f"{metric}: {value}"

                # Split the node list into chunks of 50 per slide (no truncation)
                for chunk_start in range(0, len(nodes), 50):
                    chunk_nodes = nodes[chunk_start:chunk_start + 50]

                    slide = prs.slides.add_slide(content_layout)
                    title_shape = slide.shapes.title
                    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

                    title_shape.text = category
                    if body is None:
                        continue

                    tf = body.text_frame
                    tf.clear()

                    # Main bullet for this metric
                    p_main = tf.paragraphs[0]
                    p_main.text = main_text
                    p_main.level = 0
                    _set_paragraph_font_size(p_main, MAIN_BULLET_SIZE)

                    # One level-1 bullet per node/cell in this chunk
                    for node in chunk_nodes:
                        p_node = tf.add_paragraph()
                        p_node.text = node
                        p_node.level = 1
                        _set_paragraph_font_size(p_node, SUB_BULLET_SIZE)

        # ---------------------- AUDIT CATEGORIES: single slide per category ----------------------
        else:
            slide = prs.slides.add_slide(content_layout)
            title_shape = slide.shapes.title
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

            title_shape.text = category

            if body is None:
                continue

            tf = body.text_frame
            tf.clear()

            if not items:
                p = tf.paragraphs[0]
                p.text = "No data available for this category."
                p.level = 0
                _set_paragraph_font_size(p, MAIN_BULLET_SIZE)
                continue

            if is_audit:
                # Only "Metric: Value" bullets
                for idx, item in enumerate(items):
                    metric = item.get("Metric", "")
                    value = item.get("Value", "")
                    text = f"{metric}: {value}"

                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = text
                    p.level = 0
                    _set_paragraph_font_size(p, MAIN_BULLET_SIZE)
            else:
                # Fallback: behave like audit
                for idx, item in enumerate(items):
                    metric = item.get("Metric", "")
                    value = item.get("Value", "")
                    text = f"{metric}: {value}"

                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = text
                    p.level = 0
                    _set_paragraph_font_size(p, MAIN_BULLET_SIZE)

    prs.save(ppt_path)
    return ppt_path