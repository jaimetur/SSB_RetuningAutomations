# -*- coding: utf-8 -*-

import os
from typing import Dict, List, Optional
import pandas as pd

from src.utils.utils_infrastructure import get_resource_path
from src.utils.utils_io import to_long_path


def build_text_summary_structure(
    summary_audit_df: pd.DataFrame,
) -> Dict[str, List[Dict[str, object]]]:
    """
    Group SummaryAudit rows by top-level Category, keeping SubCategory/Metric/Value/ExtraInfo.

    Returns:
      {
        "NR Frequency Audit": [
            {"SubCategory": "NRCellDU", "Metric": "...", "Value": 4, "ExtraInfo": "..."},
            ...
        ],
        "NR Frequency Inconsistencies": [
            ...
        ],
        ...
      }
    """
    sections: Dict[str, List[Dict[str, object]]] = {}

    if summary_audit_df is None or summary_audit_df.empty:
        sections["Info"] = [
            {
                "SubCategory": "Info",
                "Metric": "No audit data available to build textual summary",
                "Value": "",
                "ExtraInfo": "",
            }
        ]
        return sections

    for _, row in summary_audit_df.iterrows():
        # Group by SubCategory instead of Category
        top_key = str(row.get("SubCategory", "") or "Info")
        item = {
            "SubCategory": str(row.get("SubCategory", "") or ""),
            "Metric": str(row.get("Metric", "") or ""),
            "Value": row.get("Value", ""),
            "ExtraInfo": str(row.get("ExtraInfo", "") or ""),
        }
        sections.setdefault(top_key, []).append(item)

    return sections


def generate_ppt_summary(
    summary_audit_df: pd.DataFrame,
    excel_path: str,
    module_name: str = "",
) -> Optional[str]:
    """
    Generate a PPTX file next to the Excel with slides grouped by top-level Category.

    - First slide: global title.
    - Then, for each Category:

        • If Category name contains 'audit' (case-insensitive):
            - Single slide per Category.
            - Title = Category.
            - Body = one bullet per row: "Metric: Value" (no node list).

        • If Category name contains 'inconsist' (case-insensitive):
            - One or more slides per Category.
            - Title = Category.
            - For each row with Value > 0 and a non-empty node list:
                · Main bullet "Metric: Value".
                · Level-1 bullets with the node list parsed from ExtraInfo
                  (comma/semicolon separated), split into blocks of up to
                  100 items per slide (no truncation).

        • If Category name contains 'discrep' (case-insensitive):
            - SAME behavior as inconsistencies.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except ImportError:
        print(f"{module_name} [INFO] python-pptx is not installed. Skipping PPT summary.")
        return None

    MAIN_BULLET_SIZE = Pt(10)
    SUB_BULLET_SIZE = Pt(9)

    sections = build_text_summary_structure(summary_audit_df)

    base, _ = os.path.splitext(excel_path)
    ppt_path = base + ".pptx"
    ppt_path_long = to_long_path(ppt_path)

    def _set_paragraph_font_size(paragraph, size: Pt) -> None:
        for run in paragraph.runs:
            run.font.size = size

    def _value_is_positive(v: object) -> bool:
        """Return True if the given value represents a numeric value > 0."""
        try:
            if isinstance(v, (int, float)):
                return v > 0
            s = str(v).strip()
            if not s or s.upper() == "N/A":
                return False
            return float(s) > 0
        except Exception:
            return False

    template_path = get_resource_path("ppt_templates/ConfigurationAuditTemplate.pptx")
    try:
        prs = Presentation(template_path)
        print(f"{module_name} [INFO] Using PPT template: {template_path}")
    except Exception as e:
        print(f"{module_name} [WARNING] Could not load PPT template, using default. ({e})")
        prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # --- Title slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    title.text = "Configuration Audit Summary"
    if subtitle is not None:
        subtitle.text = os.path.basename(excel_path)

    # --- Category slides ---
    for category, items in sections.items():
        cat_lower = category.lower()
        is_audit = "audit" in cat_lower
        is_incons = "inconsist" in cat_lower  # covers 'Inconsistences' typo as well
        is_discrep = "discrep" in cat_lower  # covers 'Discrepancy/Discrepancies' and common variants
        is_incons_or_discrep = is_incons or is_discrep

        # ---------------------- INCONSISTENCIES / DISCREPANCIES: may need multiple slides ----------------------
        if is_incons_or_discrep:
            if not items:
                # No items at all: skip creating slides for this category
                continue

            for item in items:
                metric = item.get("Metric", "")
                value = item.get("Value", "")
                extra = item.get("ExtraInfo", "")

                # Skip rows with non-positive value
                if not _value_is_positive(value):
                    continue

                # Parse node/cell list from ExtraInfo
                cleaned_extra = str(extra).replace(";", ",") if extra else ""
                nodes = [t.strip() for t in cleaned_extra.split(",") if t.strip()]

                # Skip rows with empty node/cell list
                if not nodes:
                    continue

                main_text = f"{metric}: {value}"

                # Split the node list into chunks of 100 per slide (no truncation)
                for chunk_start in range(0, len(nodes), 100):
                    chunk_nodes = nodes[chunk_start:chunk_start + 100]

                    slide = prs.slides.add_slide(content_layout)
                    title_shape = slide.shapes.title
                    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

                    title_shape.text = category
                    if body is None:
                        continue

                    tf = body.text_frame
                    tf.clear()

                    # Main bullet for this metric
                    p_main = tf.paragraphs[0]
                    p_main.text = main_text
                    p_main.level = 0
                    _set_paragraph_font_size(p_main, MAIN_BULLET_SIZE)

                    # Decide how many nodes go in each column (max 4 columns × 25 nodes = 100)
                    max_per_column = 25
                    space_from_top = 0.5

                    # If there are 25 or fewer nodes, keep old single-column behavior
                    if len(chunk_nodes) <= max_per_column:
                        for node in chunk_nodes:
                            p_node = tf.add_paragraph()
                            p_node.text = node
                            p_node.level = 1
                            _set_paragraph_font_size(p_node, SUB_BULLET_SIZE)
                    else:
                        # Create chunks of up to 25 nodes per column
                        columns = [
                            chunk_nodes[i:i + max_per_column]
                            for i in range(0, len(chunk_nodes), max_per_column)
                        ]

                        # Limit to 4 columns max (100 nodes per slide)
                        columns = columns[:4]
                        num_columns = len(columns)

                        # Dynamic column width so space is equally divided
                        column_width = body.width / num_columns

                        # Create each column dynamically
                        for idx_col, col_nodes in enumerate(columns):
                            col_left = body.left + column_width * idx_col
                            col_top = body.top + Inches(space_from_top)
                            col_height = body.height

                            col_box = slide.shapes.add_textbox(col_left, col_top, column_width, col_height)
                            tf_col = col_box.text_frame
                            tf_col.clear()

                            # First node in the column
                            p_col = tf_col.paragraphs[0]
                            p_col.text = col_nodes[0]
                            p_col.level = 1
                            _set_paragraph_font_size(p_col, SUB_BULLET_SIZE)

                            for node in col_nodes[1:]:
                                p_node = tf_col.add_paragraph()
                                p_node.text = node
                                p_node.level = 1
                                _set_paragraph_font_size(p_node, SUB_BULLET_SIZE)

        # ---------------------- AUDIT CATEGORIES: single slide per category ----------------------
        else:
            slide = prs.slides.add_slide(content_layout)
            title_shape = slide.shapes.title
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

            title_shape.text = category

            if body is None:
                continue

            tf = body.text_frame
            tf.clear()

            if not items:
                p = tf.paragraphs[0]
                p.text = "No data available for this category."
                p.level = 0
                _set_paragraph_font_size(p, MAIN_BULLET_SIZE)
                continue

            if is_audit:
                # Only "Metric: Value" bullets
                for idx, item in enumerate(items):
                    metric = item.get("Metric", "")
                    value = item.get("Value", "")
                    text = f"{metric}: {value}"

                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = text
                    p.level = 0
                    _set_paragraph_font_size(p, MAIN_BULLET_SIZE)
            else:
                # Fallback: behave like audit
                for idx, item in enumerate(items):
                    metric = item.get("Metric", "")
                    value = item.get("Value", "")
                    text = f"{metric}: {value}"

                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = text
                    p.level = 0
                    _set_paragraph_font_size(p, MAIN_BULLET_SIZE)

    prs.save(ppt_path_long)
    return ppt_path
